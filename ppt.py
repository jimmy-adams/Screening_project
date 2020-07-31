# -*- coding: utf-8 -*-
"""
Created on Tue Jul 28 11:26:09 2020

@author: Melody
"""


'''
zerg 20200724 处理pptx的替换，ppt要先转换成pptx，再替换
存在问题，目录结构，只有有pptx的才会创建，如果是doc pdf 等其他文件与pptx混杂的情况可能会缺失，与原目录不一致
'''
import os
import sys
import pptx #pip install python-pptx
from pptx import Presentation
import win32com
import win32com.client
import time

TEXT_NEED_REPLACE = [('内部材料，注意保密', 'abc'), ('内部资料，注意保密', 'qianren')] #需要替换的内容，数据结构为“[(被替换内容1，替换内容1)，（被替换内容2，替换内容2），……]”

rootdir = "c:\\Users\86158\Desktop"
newrootdir = "c:\\Users\86158\Desktop"

def ppt2pptx(path):#2007以下版本先在原目录转换成pptx，删掉ppt
    for filename1 in os.listdir(path):
        subpath = os.path.join(path, filename1)
        subpath = os.path.abspath(subpath) #使用os.path.abspath('test.ppt')将相对路径转换成绝对路径
        # http://blog.sina.com.cn/s/blog_b7907b4a01014td8.html
        if subpath.endswith('.ppt'):
            print(subpath)
            try:
                #powerpoint = win32com.client.Dispatch('PowerPoint.Application')
                powerpoint = win32com.client.gencache.EnsureDispatch("PowerPoint.Application")
                powerpoint.Visible = 0
                time.sleep(2)
                ppt = powerpoint.Presentations.Open(subpath, False, False, False)
                if not os.path.exists(subpath[:-4]+'.pptx'):
                    ppt.SaveAs(subpath[:-4]+'.pptx')
                ppt.Close()
                #powerpoint.Quit()
                time.sleep(2)
                #os.remove(subpath)
                os.unlink(subpath)
            except Exception as e:
                print(os.path.split(subpath)
                      [1], "转化失败，失败原因%s" % e)


def replace_text(text_frame):#该函数实现的是文本替换功能
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            for tt in TEXT_NEED_REPLACE:
                if tt[0] in run.text:
                    run.text = run.text.replace(tt[0], tt[1])

def process_ppt(filename_open, filename_save):
    prs = Presentation(filename_open)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:#判断Shape是否含有文本框
                text_frame = shape.text_frame
                replace_text(text_frame)#调用replace_text函数实现文本替换
            if shape.has_table:#判断Shape是否含有表格
                table = shape.table
                for cell in table.iter_cells():#遍历表格的cell
                    text_frame = cell.text_frame
                    replace_text(text_frame)#调用replace_text函数实现文本替换
    newpath = os.path.dirname(filename_save)
    if not os.path.exists(newpath):
        os.mkdir(newpath)
    prs.save(filename_save) #保存

orapath = []

dirs = os.listdir(rootdir)#列出文件夹下所有的目录与文件
g = os.walk(rootdir)
for path, dir_list, file_list in g:#老版本要转成pptx
    for dir_name in dir_list:
        ppt2pptx(os.path.join(path, dir_name))

dirs = os.listdir(rootdir)#列出文件夹下所有的目录与文件
g = os.walk(rootdir)
for path, dir_list, file_list in g:
    for file_name in file_list:
        if file_name.endswith(".pptx"):
            filepath = os.path.abspath(os.path.join(path, file_name))
            print(filepath)
            orapath.append(filepath)
            FILE_OPEN = filepath
            FILE_SAVE = filepath.replace(rootdir, newrootdir)

            process_ppt(FILE_OPEN, FILE_SAVE)
#    print(os.path.join(path, file))


print(len(orapath))