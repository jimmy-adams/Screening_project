# -*- coding: utf-8 -*-
"""
Created on Thu Jul 23 16:44:41 2020

@author: Melody
"""
import os
from docx import Document #pip install python-docx
import pptx #pip install python-pptx
from pptx import Presentation
from win32com import client as wc
# 放了一些docx 文件
old_file_path = "C:/Users/86158/Desktop/"
# 生成新文件后的存放地址
new_file_path = "C:/Users/86158/Desktop/"

replace_dict = {
    "深圳":"ShenZhen",
    "Modified":"",
}



def check_and_change(document, replace_dict, id_type):
    """
    遍历word中的所有 paragraphs，在每一段中发现含有key 的内容，就替换为 value 。 
    （key 和 value 都是replace_dict中的键值对。）
    """
    if id_type == 1:
        #=======Process on the paragraph ========================================#
        for para in document.paragraphs:
            for i in range(len(para.runs)):
                for key, value in replace_dict.items():
                    if key in para.runs[i].text:
                        print(key+"-->"+value)
                        para.runs[i].text = para.runs[i].text.replace(key, value)  
       #======Process on the table==============================================#
        for table in document.tables[:]:
           for i, row in enumerate(table.rows[:]):   # 读每行
               for cell in row.cells[:]:  # 读一行中的所有单元格
                   for key, value in replace_dict.items():
                        if key in cell.text:
                            print(key+"-->"+value)
                            cell.text = cell.text.replace(key,value)
       #========================================================================#
       #======Process on the header & foot======================================#
        for section in document.sections:
           header = section.header 
           foot   = section.footer
           for key, value in replace_dict.items():
               if key in header.paragraphs[0].text:
                   header.paragraphs[0].text = header.paragraphs[0].text.replace(key,value)
                   print(header.paragraphs[0].text)
               if key in foot.paragraphs[0].text:
                   foot.paragraphs[0].text = foot.paragraphs[0].text.replace(key,value)
                   print(foot.paragraphs[0].text)
       #========================================================================#
    elif id_type == 2:
        #======Process on PPTX===================================================#
        for slide in document.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:#判断Shape是否含有文本框
                    text_frame = shape.text_frame
                    for para in text_frame.paragraphs:
                        for i in range(len(para.runs)):
                            for key, value in replace_dict.items():
                                if key in para.runs[i].text:
                                    print(key+"-->"+value)
                                    para.runs[i].text = para.runs[i].text.replace(key, value)
                if shape.has_table:#判断Shape是否含有表格
                    table = shape.table
                    for cell in table.iter_cells():#遍历表格的cell
                        text_frame = cell.text_frame
                        for para in text_frame.paragraphs:
                            for i in range(len(para.runs)):
                                for key, value in replace_dict.items():
                                     if key in para.runs[i].text:
                                         print(key+"-->"+value)
                                         para.runs[i].text = para.runs[i].text.replace(key, value)
       #========================================================================#            
    return document   
    

def main():
    word = wc.Dispatch("Word.Application")
    powerpoint = wc.Dispatch("PowerPoint.Application")
    for name in os.listdir(old_file_path):
        print(name)
        old_file = old_file_path + name
        new_file = old_file_path + name
        
        #==========preprocess for doc===============================================#
        #==========convert the doc to docx,and delete the doc file==================#
        if old_file.split(".")[-1] == 'doc':
            doc = word.Documents.Open(old_file)
            doc.SaveAs("{}x".format(old_file), 12)#另存为后缀为".docx"的文件，其中参数12指docx文件
            doc.Close() #关闭原来word文件
            os.remove(old_file)
        #============================================================================#
        if old_file.split(".")[-1] == 'docx':
            id_type = 1
            document = Document(old_file)
            document = check_and_change(document, replace_dict,id_type)
            document.save(new_file)
        #==========convert the ppt to pptx,and delete the ppt file==================#
        if old_file.split(".")[-1] == 'ppt':
            ppt = powerpoint.Presentations.Open(old_file)
            ppt.SaveAs(old_file[:-4]+'.pptx')#另存为后缀为".pptx"的文件，其中参数12指pptx文件
            os.remove(old_file)
            ppt.Close() #关闭原来word文件
            #os.unlink(old_file)
        #============================================================================#
        if old_file.split(".")[-1] == 'pptx':
            id_type = 2
            document = Presentation(old_file)
            document = check_and_change(document, replace_dict,id_type)
            document.save(new_file)
        
        print("="*30)
        

if __name__ == '__main__':
    main()

